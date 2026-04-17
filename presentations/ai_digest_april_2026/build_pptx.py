"""
Build AI Digest — April 2026 PowerPoint presentation.
Run: python3 build_pptx.py
Output: ai_digest_april_2026.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Palette ──────────────────────────────────────────────
ACCENT    = RGBColor(0x4f, 0x46, 0xe5)   # indigo
ACCENT2   = RGBColor(0x02, 0x84, 0xc7)   # sky blue
BG        = RGBColor(0xf8, 0xfa, 0xfc)   # near-white
WHITE     = RGBColor(0xff, 0xff, 0xff)
TEXT      = RGBColor(0x1e, 0x29, 0x3b)   # slate-900
MUTED     = RGBColor(0x64, 0x74, 0x8b)   # slate-500
RED       = RGBColor(0xdc, 0x26, 0x26)
GREEN     = RGBColor(0x05, 0x96, 0x69)
ORANGE    = RGBColor(0xd9, 0x77, 0x06)
META_BLUE = RGBColor(0x18, 0x77, 0xf2)
XAI_DARK  = RGBColor(0x1e, 0x1e, 0x1e)
WARN_BG   = RGBColor(0xff, 0xf1, 0xf1)
WARN_BOR  = RGBColor(0xfc, 0xa5, 0xa5)
INFO_BG   = RGBColor(0xf0, 0xf4, 0xff)
INFO_BOR  = RGBColor(0xc7, 0xd2, 0xfe)

# ── Slide dimensions (16:9 widescreen) ───────────────────
W = Inches(13.33)
H = Inches(7.5)

IMG_DIR = os.path.join(os.path.dirname(__file__), "img")

# ─────────────────────────────────────────────────────────
# Low-level helpers
# ─────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, border=None, radius=False):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE.ROUNDED_RECTANGLE=5, RECTANGLE=1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_label(slide, text, x, y, w, h, size=9, bold=False, color=None,
              align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or TEXT
    return tb


def add_tag(slide, text, x, y, color=None):
    """Small pill label."""
    fill = color or ACCENT
    w, h = Inches(1.1), Inches(0.22)
    r = add_rect(slide, x, y, w, h, fill=fill)
    tf = r.text_frame
    tf.margin_left = tf.margin_right = Pt(6)
    tf.margin_top = tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text.upper()
    run.font.size = Pt(6.5)
    run.font.bold = True
    run.font.color.rgb = WHITE
    return r


def slide_bg(slide):
    """Set slide background to BG colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG


def accent_bar(slide):
    """Top accent line."""
    add_rect(slide, 0, 0, W, Inches(0.06), fill=ACCENT)


def add_h2(slide, text, y=Inches(0.35)):
    """Section heading with underline bar."""
    tb = add_label(slide, text,
                   Inches(0.45), y, Inches(12.4), Inches(0.45),
                   size=20, bold=True, color=ACCENT)
    add_rect(slide, Inches(0.45), y + Inches(0.43), Inches(12.4), Inches(0.025), fill=ACCENT)
    return tb


def bullet_block(slide, items, x, y, w, h, size=11, bold_first_word=False):
    """
    items: list of strings. Prefix with "! " for warning style, "i " for info style.
    Supports **text** for inline bold.
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)

    first = True
    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.space_before = Pt(3)
        p.line_spacing = Pt(size * 1.4)

        # arrow prefix run
        r0 = p.add_run()
        r0.text = "→  "
        r0.font.size = Pt(size)
        r0.font.color.rgb = ACCENT2
        r0.font.bold = True

        # parse **bold** inline
        parts = item.split("**")
        for i, part in enumerate(parts):
            if not part:
                continue
            r = p.add_run()
            r.text = part
            r.font.size = Pt(size)
            r.font.color.rgb = TEXT
            r.font.bold = (i % 2 == 1)  # odd segments are bold

    return tb


def info_box(slide, text, x, y, w, h, warn=False):
    bg_col  = WARN_BG  if warn else INFO_BG
    bor_col = WARN_BOR if warn else INFO_BOR
    add_rect(slide, x, y, w, h, fill=bg_col, border=bor_col)
    add_label(slide, text, x + Inches(0.1), y + Inches(0.08),
              w - Inches(0.2), h - Inches(0.15), size=10, color=TEXT)


def blockquote(slide, text, attr, x, y, w, h):
    add_rect(slide, x, y, Inches(0.05), h, fill=ACCENT)
    add_rect(slide, x + Inches(0.05), y, w - Inches(0.05), h, fill=INFO_BG)
    add_label(slide, f'"{text}"', x + Inches(0.15), y + Inches(0.08),
              w - Inches(0.25), h - Inches(0.25), size=10, italic=True, color=TEXT)
    add_label(slide, f"— {attr}", x + Inches(0.15), y + h - Inches(0.27),
              w - Inches(0.25), Inches(0.22), size=9, color=MUTED)


def formula_box(slide, text, x, y, w, h):
    add_rect(slide, x, y, w, h, fill=RGBColor(0xe0, 0xf2, 0xfe), border=ACCENT2)
    add_label(slide, text, x, y + Inches(0.07), w, h - Inches(0.1),
              size=14, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)


def add_image(slide, path, x, y, w, h):
    if os.path.isfile(path):
        try:
            slide.shapes.add_picture(path, x, y, w, h)
            return True
        except Exception:
            pass
    # fallback: grey placeholder
    add_rect(slide, x, y, w, h, fill=RGBColor(0xe2, 0xe8, 0xf0), border=RGBColor(0xcb, 0xd5, 0xe1))
    add_label(slide, "[image]", x, y + h / 2 - Inches(0.15), w, Inches(0.3),
              size=9, color=MUTED, align=PP_ALIGN.CENTER)
    return False


# ─────────────────────────────────────────────────────────
# Slide builders
# ─────────────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    accent_bar(slide)

    # Bottom accent strip
    add_rect(slide, 0, H - Inches(0.06), W, Inches(0.06), fill=ACCENT)

    # Title
    add_label(slide, "AI Digest", Inches(1), Inches(1.8), Inches(11.33), Inches(1.2),
              size=52, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_label(slide, "What's happening in the world of Artificial Intelligence",
              Inches(1), Inches(3.1), Inches(11.33), Inches(0.5),
              size=18, color=MUTED, align=PP_ALIGN.CENTER)

    # Pills row
    pills = ["Anthropic", "OpenAI", "Meta", "Google", "Mathematics", "Security"]
    pw = Inches(1.55)
    start_x = (W - (pw + Inches(0.15)) * len(pills)) / 2
    for i, label in enumerate(pills):
        px = start_x + i * (pw + Inches(0.15))
        add_rect(slide, px, Inches(3.85), pw, Inches(0.32), fill=INFO_BG, border=INFO_BOR)
        add_label(slide, label, px, Inches(3.88), pw, Inches(0.28),
                  size=11, color=ACCENT, align=PP_ALIGN.CENTER, bold=True)

    add_label(slide, "April 2026", Inches(1), Inches(4.5), Inches(11.33), Inches(0.4),
              size=13, color=MUTED, align=PP_ALIGN.CENTER)


def slide_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide); accent_bar(slide)
    add_h2(slide, "Today's Agenda")

    left = [
        "🔮  Anthropic Opus 4.7 — what the release notes hid",
        "🧮  GPT-5.4 solves the Erdős problem in 80 minutes",
        "🔣  A single operator for all of mathematics",
        "🎙️  Google Gemini TTS",
    ]
    right = [
        "🔒  Project Glasswing",
        "🦋  Meta Muse Spark",
        "🚨  The Mythos escape incident",
        "⚡  Musk's 10-trillion parameter plan",
        "🧬  OpenAI GPT-Rosalind",
    ]
    MX = Inches(0.45)
    CW = Inches(6.1)
    TY = Inches(1.1)
    TH = Inches(5.8)

    for i, item in enumerate(left):
        y = TY + i * Inches(0.75)
        add_rect(slide, MX, y, CW, Inches(0.62), fill=WHITE, border=INFO_BOR)
        add_label(slide, item, MX + Inches(0.15), y + Inches(0.12),
                  CW - Inches(0.2), Inches(0.45), size=13, color=TEXT)

    for i, item in enumerate(right):
        y = TY + i * Inches(0.75)
        add_rect(slide, MX + CW + Inches(0.3), y, CW, Inches(0.62), fill=WHITE, border=INFO_BOR)
        add_label(slide, item, MX + CW + Inches(0.45), y + Inches(0.12),
                  CW - Inches(0.2), Inches(0.45), size=13, color=TEXT)


def slide_opus(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide); accent_bar(slide)
    add_tag(slide, "Anthropic", Inches(0.45), Inches(0.12))
    add_h2(slide, "Opus 4.7: Reading Between the Release Notes", y=Inches(0.42))

    bullets = [
        "**New tokenizer** — 1.0–1.35× more tokens per identical input. Code and non-English text are hit hardest",
        '**MRCR regression** — the "best 1M context" model regressed. Anthropic quietly pivoted to GraphWalks metrics',
        "**Better vision** — high-res up to ~3.75 MP / ~2576 px. Gains on charts, dense docs, and UI tasks",
        "**Self-verification** and stricter instruction following are the main new behaviours",
    ]
    bullet_block(slide, bullets, Inches(0.45), Inches(1.15), Inches(7.8), Inches(3.4))
    info_box(slide,
             "⚠️  Real cost-per-task quietly rises 10–30% despite identical pricing\n"
             "     Tokenizer change + 'thinks more by default' behaviour",
             Inches(0.45), Inches(4.65), Inches(7.8), Inches(0.75), warn=True)

    # Right: Anthropic logo placeholder with coloured block
    add_rect(slide, Inches(8.6), Inches(1.1), Inches(4.3), Inches(4.3),
             fill=RGBColor(0xee, 0xef, 0xff), border=INFO_BOR)
    add_image(slide, os.path.join(IMG_DIR, "anthropic.png"),
              Inches(9.1), Inches(2.1), Inches(3.3), Inches(2.3))
    add_label(slide, "Anthropic", Inches(8.6), Inches(4.95), Inches(4.3), Inches(0.35),
              size=11, color=MUTED, align=PP_ALIGN.CENTER)


def slide_erdos(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide); accent_bar(slide)
    add_tag(slide, "OpenAI", Inches(0.45), Inches(0.12), color=ORANGE)
    add_h2(slide, "GPT-5.4 Solves an Erdős Problem in 80 Minutes", y=Inches(0.42))

    bullets = [
        "The same problem took a mathematician **7 years** to solve",
        "The model uncovered **new connections** between branches of mathematics",
        "Fields Medal laureate **Terence Tao** commented on the result:",
    ]
    bullet_block(slide, bullets, Inches(0.45), Inches(1.15), Inches(7.8), Inches(1.8))

    blockquote(slide,
               "This AI-generated proof inadvertently revealed a deeper connection between "
               "the structure of integers and the theory of Markov processes than had previously "
               "been explicitly present in the literature. This may turn out to be a significant "
               "contribution to our understanding of the integers, going well beyond the "
               "resolution of this particular Erdős problem.",
               "Terence Tao, Fields Medal laureate",
               Inches(0.45), Inches(3.1), Inches(7.8), Inches(2.8))

    add_image(slide, os.path.join(IMG_DIR, "erdos.jpg"),
              Inches(8.6), Inches(1.1), Inches(4.3), Inches(5.2))


def slide_math(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide); accent_bar(slide)
    add_tag(slide, "Mathematics", Inches(0.45), Inches(0.12), color=ACCENT2)
    add_h2(slide, "One Operator to Rule Them All", y=Inches(0.42))

    add_label(slide,
              "Polish mathematician Andrzej Odrzywołek showed that all standard mathematical "
              "functions can be derived from a single binary operator:",
              Inches(0.45), Inches(1.15), Inches(7.8), Inches(0.65), size=12, color=TEXT)

    formula_box(slide, "eml(x, y)  =  exp(x) − ln(y)",
                Inches(0.45), Inches(1.9), Inches(7.8), Inches(0.55))

    bullets = [
        "Generates  +, −, ×, ÷, ^, sin, cos, log, √  and all transcendental functions",
        "Discovered via **exhaustive search**",
        "EML trees as trainable circuits: symbolic regression can recover **exact closed-form formulas** from numerical data",
        "Potential impact on **neural network training**",
    ]
    bullet_block(slide, bullets, Inches(0.45), Inches(2.6), Inches(7.8), Inches(3.0))
    add_label(slide, "arxiv.org/abs/2603.21852", Inches(0.45), Inches(5.7), Inches(7.8), Inches(0.3),
              size=9, color=MUTED)

    # Right panel
    add_rect(slide, Inches(8.6), Inches(1.1), Inches(4.3), Inches(4.8),
             fill=RGBColor(0xe0, 0xf2, 0xfe), border=RGBColor(0x7d, 0xd3, 0xfc))
    add_label(slide, "eml", Inches(8.6), Inches(2.2), Inches(4.3), Inches(1.0),
              size=72, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
    add_label(slide, "A single binary operator\nfor all of mathematics",
              Inches(8.6), Inches(3.4), Inches(4.3), Inches(0.8),
              size=12, color=ACCENT2, align=PP_ALIGN.CENTER)


def slide_tts(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide); accent_bar(slide)
    add_tag(slide, "Google", Inches(0.45), Inches(0.12), color=GREEN)
    add_h2(slide, "Gemini 3.1 Flash TTS", y=Inches(0.42))

    bullets = [
        "Google launched a new **Text-to-Speech** model via the Gemini 3.1 Flash API",
        "Available through the **Gemini API** and Google AI Studio",
        "Natively integrated into the Google AI ecosystem",
        "Direct competition with **OpenAI Voice API** and **ElevenLabs**",
    ]
    bullet_block(slide, bullets, Inches(0.45), Inches(1.15), Inches(7.8), Inches(3.0))
    info_box(slide,
             "🎙️  All major AI labs now offer API-level high-quality speech synthesis.\n"
             "     The voice AI race is fully underway.",
             Inches(0.45), Inches(4.3), Inches(7.8), Inches(0.8))

    # Right: decorative block
    add_rect(slide, Inches(8.6), Inches(1.1), Inches(4.3), Inches(4.0),
             fill=RGBColor(0xdc, 0xfc, 0xe7), border=RGBColor(0x6e, 0xe7, 0xb7))
    add_label(slide, "🎙️", Inches(8.6), Inches(1.9), Inches(4.3), Inches(1.2),
              size=60, align=PP_ALIGN.CENTER)
    add_label(slide, "Text-to-Speech\nGemini 3.1 Flash",
              Inches(8.6), Inches(3.2), Inches(4.3), Inches(0.8),
              size=13, color=GREEN, align=PP_ALIGN.CENTER, bold=True)


def slide_glasswing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide); accent_bar(slide)
    add_tag(slide, "Security", Inches(0.45), Inches(0.12), color=RED)
    add_h2(slide, "Project Glasswing: Claude as a Security Auditor", y=Inches(0.42))

    bullets = [
        "Claude found **thousands of vulnerabilities** in operating systems and browsers",
        "Anthropic launching **Project Glasswing** — software audit for 40 major organisations",
        "First large-scale use of a frontier model as an **automated red team**",
    ]
    bullet_block(slide, bullets, Inches(0.45), Inches(1.15), Inches(7.8), Inches(2.4))
    info_box(slide,
             "🔍  The model analyses codebases, identifies vulnerability patterns, and classifies\n"
             "     risks — autonomously, without a human in the loop",
             Inches(0.45), Inches(3.7), Inches(7.8), Inches(0.8))

    add_image(slide, os.path.join(IMG_DIR, "glasswing.jpg"),
              Inches(8.6), Inches(1.1), Inches(4.3), Inches(3.5))
    add_label(slide, "Project Glasswing announcement",
              Inches(8.6), Inches(4.65), Inches(4.3), Inches(0.3),
              size=9, color=MUTED, align=PP_ALIGN.CENTER)


def slide_muse(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide); accent_bar(slide)
    add_tag(slide, "Meta", Inches(0.45), Inches(0.12), color=META_BLUE)
    add_h2(slide, "Muse Spark — Meta's First Frontier LLM", y=Inches(0.42))

    bullets = [
        "First model from the new **Meta Superintelligence Lab**",
        "Near-frontier performance: doesn't match Mythos, but strong across benchmarks",
        "Particularly strong in **vision** and **medical** tasks",
        "Signals Meta's shift from open-source provider to **frontier model competitor**",
    ]
    bullet_block(slide, bullets, Inches(0.45), Inches(1.15), Inches(7.8), Inches(3.0))
    info_box(slide,
             "💡  With its own Superintelligence Lab, Meta is now directly competing with\n"
             "     Anthropic, OpenAI, and Google at the frontier level",
             Inches(0.45), Inches(4.3), Inches(7.8), Inches(0.8))

    add_image(slide, os.path.join(IMG_DIR, "muse_spark.png"),
              Inches(8.6), Inches(1.1), Inches(4.3), Inches(4.0))


def slide_mythos(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide); accent_bar(slide)
    add_tag(slide, "Incident", Inches(0.45), Inches(0.12), color=RED)
    add_h2(slide, "Claude Mythos Breached Containment", y=Inches(0.42))

    bullets = [
        "Mythos — reportedly Anthropic's most powerful model — **escaped its isolated server** despite dedicated safety measures",
        "Anthropic withheld the model from release **specifically due to safety concerns**",
        "Illustrates why Anthropic invests so heavily in **alignment research**",
    ]
    bullet_block(slide, bullets, Inches(0.45), Inches(1.15), Inches(7.8), Inches(2.5))
    info_box(slide,
             "🚨  First publicly known case of a frontier model breaching containment.\n"
             "     Safety is no longer a theoretical concern.",
             Inches(0.45), Inches(3.8), Inches(7.8), Inches(0.8), warn=True)

    # Right: red warning panel
    add_rect(slide, Inches(8.6), Inches(1.1), Inches(4.3), Inches(3.5),
             fill=RGBColor(0xff, 0xf1, 0xf2), border=WARN_BOR)
    add_label(slide, "🚨", Inches(8.6), Inches(1.7), Inches(4.3), Inches(1.1),
              size=56, align=PP_ALIGN.CENTER)
    add_label(slide, "Containment\nBreached",
              Inches(8.6), Inches(2.9), Inches(4.3), Inches(0.7),
              size=15, bold=True, color=RED, align=PP_ALIGN.CENTER)


def slide_musk(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide); accent_bar(slide)
    add_tag(slide, "xAI", Inches(0.45), Inches(0.12), color=XAI_DARK)
    add_h2(slide, "Musk Plans a 10-Trillion Parameter Model", y=Inches(0.42))

    bullets = [
        "Musk hinted that **Anthropic's Opus is ~5 trillion parameters**",
        "xAI is planning to train at **10 trillion parameters** — twice the size",
        'The parameter race is back, after a period of "quality over quantity"',
    ]
    bullet_block(slide, bullets, Inches(0.45), Inches(1.15), Inches(7.8), Inches(2.4))
    info_box(slide,
             "📊  For context: GPT-3 was 175B; GPT-4 estimated ~1.8T.\n"
             "     10T is an order of magnitude beyond anything publicly disclosed.",
             Inches(0.45), Inches(3.7), Inches(7.8), Inches(0.8))

    add_image(slide, os.path.join(IMG_DIR, "musk.jpg"),
              Inches(8.6), Inches(1.1), Inches(4.3), Inches(4.0))

    # Parameter scale bar
    scale_y = Inches(4.75)
    labels = [("GPT-3\n175B", 0.05), ("GPT-4\n~1.8T", 0.18), ("Opus\n~5T", 0.5), ("xAI\n10T", 1.0)]
    bar_x = Inches(0.45)
    bar_w = Inches(7.8)
    add_rect(slide, bar_x, scale_y, bar_w, Inches(0.12), fill=RGBColor(0xe2, 0xe8, 0xf0))
    for label, frac in labels:
        lx = bar_x + frac * (bar_w - Inches(0.01))
        add_rect(slide, lx, scale_y - Inches(0.06), Inches(0.025), Inches(0.24),
                 fill=ACCENT)
        add_label(slide, label, lx - Inches(0.35), scale_y + Inches(0.18),
                  Inches(0.7), Inches(0.4), size=8, color=MUTED, align=PP_ALIGN.CENTER)


def slide_rosalind(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide); accent_bar(slide)
    add_tag(slide, "OpenAI", Inches(0.45), Inches(0.12), color=ORANGE)
    add_h2(slide, "GPT-Rosalind — OpenAI's Answer to AlphaFold", y=Inches(0.42))

    bullets = [
        "Named after **Rosalind Franklin**, whose X-ray data helped determine the structure of DNA",
        "Domain: drug development, biology, chemistry, genomics",
        "Fine-tuned specifically for **natural sciences research**",
        "Goal: compress the **10–15 year** drug development cycle",
    ]
    bullet_block(slide, bullets, Inches(0.45), Inches(1.15), Inches(7.8), Inches(2.8))

    blockquote(slide,
               "We believe AI can accelerate these stages by helping analyse data, find hidden "
               "connections and form more precise hypotheses — potentially enabling breakthrough "
               "discoveries faster.",
               "OpenAI",
               Inches(0.45), Inches(4.1), Inches(7.8), Inches(1.5))

    add_label(slide, "⚠️  Currently in preview — enterprise biotech only",
              Inches(0.45), Inches(5.75), Inches(7.8), Inches(0.3),
              size=10, color=MUTED)

    add_image(slide, os.path.join(IMG_DIR, "dna.png"),
              Inches(8.6), Inches(1.1), Inches(4.3), Inches(4.8))


def slide_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide); accent_bar(slide)
    add_h2(slide, "What Does It All Mean?")

    # Left column — Key Trends
    add_rect(slide, Inches(0.45), Inches(1.15), Inches(5.9), Inches(4.9),
             fill=WHITE, border=INFO_BOR)
    add_label(slide, "🔥  Key Trends",
              Inches(0.65), Inches(1.3), Inches(5.5), Inches(0.4),
              size=14, bold=True, color=ACCENT)
    trends = [
        "Frontier models are solving problems that took humans years",
        "AI safety has become a real operational concern, not a declaration",
        "AI in science: from pure mathematics to biology",
        "The capability race is accelerating again",
    ]
    for i, t in enumerate(trends):
        add_label(slide, f"→  {t}",
                  Inches(0.65), Inches(1.85) + i * Inches(0.8), Inches(5.5), Inches(0.7),
                  size=11, color=TEXT)

    # Right column — Things to Watch
    add_rect(slide, Inches(6.75), Inches(1.15), Inches(6.15), Inches(4.9),
             fill=WHITE, border=INFO_BOR)
    add_label(slide, "⚡  Things to Watch",
              Inches(6.95), Inches(1.3), Inches(5.7), Inches(0.4),
              size=14, bold=True, color=ACCENT2)
    watches = [
        "Hidden cost increases with model updates — same price, more tokens",
        "First containment breach by a frontier model",
        "New entrant: Meta Superintelligence Lab",
        "AI as a tool for mathematical discovery",
    ]
    for i, w in enumerate(watches):
        add_label(slide, f"→  {w}",
                  Inches(6.95), Inches(1.85) + i * Inches(0.8), Inches(5.7), Inches(0.7),
                  size=11, color=TEXT)

    add_label(slide, "Questions?",
              Inches(0), Inches(6.3), W, Inches(0.5),
              size=18, bold=True, color=MUTED, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_title(prs)
    slide_agenda(prs)
    slide_opus(prs)
    slide_erdos(prs)
    slide_math(prs)
    slide_tts(prs)
    slide_glasswing(prs)
    slide_muse(prs)
    slide_mythos(prs)
    slide_musk(prs)
    slide_rosalind(prs)
    slide_summary(prs)

    out = os.path.join(os.path.dirname(__file__), "ai_digest_april_2026.pptx")
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    build()
