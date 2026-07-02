"""One slide: GRPO vs our final method (GTPO-EMA-flipped FIXED, λ=0.5, pos_discount,
top-k=5) — reward/advantage math + training process, side by side."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

HERE = os.path.dirname(__file__)
NAVY = RGBColor(0x1e, 0x3a, 0x5f); GREY = RGBColor(0x64, 0x74, 0x8b)
GREEN = RGBColor(0x15, 0x80, 0x3d); DARK = RGBColor(0x22, 0x22, 0x22)
LGREY = RGBColor(0xf1, 0xf5, 0xf9); LGREEN = RGBColor(0xe7, 0xf5, 0xec)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])


def box(x, y, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill if fill else RGBColor(0xff, 0xff, 0xff)
    if line:
        sh.line.color.rgb = line; sh.line.width = Pt(1.5)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def txt(x, y, w, h, lines, size=11, color=DARK, bold_first=False, align=PP_ALIGN.LEFT, mono=False):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf = tb.text_frame
    tf.word_wrap = True
    for i, (t, opts) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", align); p.space_after = Pt(opts.get("sa", 3)); p.space_before = Pt(opts.get("sb", 0))
        r = p.add_run(); r.text = t
        f = r.font; f.size = Pt(opts.get("size", size)); f.bold = opts.get("bold", False)
        f.color.rgb = opts.get("color", color)
        f.name = "Consolas" if opts.get("mono", mono) else "Calibri"
    return tb


# Title
txt(0.4, 0.15, 12.5, 0.7, [
    ("GRPO  vs  GTPO-EMA-flipped (FIXED) + pos_discount   —   reward & training", {"size": 22, "bold": True, "color": NAVY}),
    ("same rollouts, rewards, PPO-clip loss;  differ ONLY in the token-level advantage Ã", {"size": 12, "color": GREY}),
])

COLY, COLH = 1.15, 5.7
# ---- LEFT: GRPO ----
box(0.35, COLY, 6.15, COLH, fill=LGREY)
txt(0.55, COLY + 0.1, 5.8, 0.4, [("GRPO (baseline)", {"size": 16, "bold": True, "color": GREY})])
txt(0.55, COLY + 0.65, 5.75, 5.0, [
    ("Rollouts", {"size": 12, "bold": True, "color": NAVY}),
    ("group of G completions {o_i} per prompt q;", {"size": 11}),
    ("terminal reward r_i  (correct / wrong).", {"size": 11}),
    ("Advantage  (uniform per sequence)", {"size": 12, "bold": True, "color": NAVY, "sb": 8}),
    ("Â_i = (r_i − mean(r)) / std(r)", {"size": 13, "mono": True, "color": DARK}),
    ("→ the SAME Â_i on every token t of o_i.", {"size": 11}),
    ("Objective  (token-mean PPO-clip)", {"size": 12, "bold": True, "color": NAVY, "sb": 8}),
    ("J = E[ (1/G)Σ_i (1/|o_i|)Σ_t", {"size": 11.5, "mono": True}),
    ("      min( w_it·Â_i , clip(w_it,1±ε)·Â_i ) ]", {"size": 11.5, "mono": True}),
    ("w_it = π_θ(o_it)/π_θold(o_it)", {"size": 11, "mono": True, "color": GREY}),
    ("Credit assignment", {"size": 12, "bold": True, "color": NAVY, "sb": 8}),
    ("COARSE — one scalar shared by all tokens;", {"size": 11}),
    ("no notion of which tokens matter.", {"size": 11}),
])

# ---- RIGHT: OURS ----
box(6.75, COLY, 6.25, COLH, fill=LGREEN, line=GREEN)
txt(6.95, COLY + 0.1, 5.9, 0.5, [
    ("Ours: GTPO-EMA-flipped (FIXED)", {"size": 16, "bold": True, "color": GREEN}),
    ("λ = 0.5 · pos_discount · top-k = 5", {"size": 11.5, "bold": True, "color": GREEN}),
])
txt(6.95, COLY + 0.85, 5.9, 4.8, [
    ("Per-token confidence & EMA", {"size": 12, "bold": True, "color": NAVY}),
    ("C_it = −(1/k) Σ_{v∈top-k} log π(v),  k=5", {"size": 11.5, "mono": True}),
    ("EMA_it = λ·EMA_i,t−1 + (1−λ)·C_it,  λ=0.5", {"size": 11.5, "mono": True}),
    ("Split by terminal reward: O⁺ (correct) / O⁻ (wrong)", {"size": 11, "sb": 2}),
    ("Shaped per-token reward  (flipped roles)", {"size": 12, "bold": True, "color": NAVY, "sb": 7}),
    ("O⁺: r̃_it = α₁ + α₂·g(t)·(EMA_it⁻¹ / Σ_{O⁺}EMA⁻¹)·d_t", {"size": 11, "mono": True, "color": GREEN}),
    ("O⁻: r̃_jt = −[ α₁ + α₂·g(t)·(EMA_jt / Σ_{O⁻}EMA)·h_t ]", {"size": 11, "mono": True, "color": RGBColor(0xb9,0x1c,0x1c)}),
    ("pos_discount:  g(t) = τ/(τ+t),  τ=1024", {"size": 11, "mono": True, "color": GREY}),
    ("(reward low-confidence/exploratory tokens on correct", {"size": 10, "color": GREY, "sa": 0}),
    (" paths; punish confident tokens on wrong ones; damp late)", {"size": 10, "color": GREY}),
    ("Per-token advantage", {"size": 12, "bold": True, "color": NAVY, "sb": 6}),
    ("Ã_it = z_{O⁺}(r̃⁺) + z_{O⁻}(r̃⁻)   (per-polarity z-norm)", {"size": 11, "mono": True}),
    ("FIXED: computed on the FULL group at generation", {"size": 11, "sb": 4}),
    ("(π_θold), injected per-token — NOT the degenerate B=1", {"size": 10.5, "color": GREY, "sa": 0}),
    (" recompute that made length explode.", {"size": 10.5, "color": GREY}),
    ("Objective: same PPO-clip, with Â_i → Ã_it (per token)", {"size": 11, "bold": True, "color": NAVY, "sb": 5}),
])

# footer
txt(0.4, 6.98, 12.5, 0.4, [
    ("Empirical (Qwen3-4B-Base): fine-grained credit + top-k=5 confidence beats GRPO on GSM8K / MATH-500 / Big-Math at shorter length; hard (Omni-MATH) unchanged.",
     {"size": 10.5, "color": GREY, "align": PP_ALIGN.CENTER})])

out = os.path.join(HERE, "final_method_vs_grpo.pptx")
prs.save(out)
print("saved", out)
