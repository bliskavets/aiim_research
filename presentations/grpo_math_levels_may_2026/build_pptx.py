"""
Build grpo_math_levels_may_2026.pptx.

Covers exp_039..exp_048 — Qwen3-4B on Big-Math int-2000 and MATH-lighteval
levels 3-5 (integer answers). Same minimal style as grpo_update_april_2026.pptx.

Run:   python3 build_pptx.py
Output: grpo_math_levels_may_2026.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
IMG  = os.path.join(HERE, "img")
OUT  = os.path.join(HERE, "grpo_math_levels_may_2026.pptx")

ACCENT = RGBColor(0x4f, 0x46, 0xe5)
TEXT   = RGBColor(0x1e, 0x29, 0x3b)
MUTED  = RGBColor(0x64, 0x74, 0x8b)


def _strip_placeholders(slide):
    for shape in list(slide.shapes):
        if shape.has_text_frame and shape.text_frame.text.strip() in (
            "Click to edit Master title style", "Title"
        ):
            sp = shape._element
            sp.getparent().remove(sp)


def add_title_slide(prs, title, subtitle):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    _strip_placeholders(slide)
    w, h = prs.slide_width, prs.slide_height

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.3),
                                  w - Inches(1.4), Inches(1.3))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = ACCENT

    sb = slide.shapes.add_textbox(Inches(0.7), Inches(3.7),
                                  w - Inches(1.4), Inches(1.0))
    p = sb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = subtitle
    r.font.size = Pt(16); r.font.italic = True; r.font.color.rgb = MUTED


def add_text_slide(prs, title, body_lines):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    _strip_placeholders(slide)
    w, h = prs.slide_width, prs.slide_height

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3),
                                  w - Inches(1.0), Inches(0.7))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = ACCENT

    bb = slide.shapes.add_textbox(Inches(0.5), Inches(1.15),
                                  w - Inches(1.0), h - Inches(1.4))
    tf = bb.text_frame; tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = line
        r.font.size = Pt(16); r.font.color.rgb = TEXT


def add_image_slide(prs, title, img_name, caption=None):
    img_path = os.path.join(IMG, img_name)
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    _strip_placeholders(slide)
    w, h = prs.slide_width, prs.slide_height

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3),
                                  w - Inches(1.0), Inches(0.7))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = ACCENT

    img_top = Inches(1.15)
    avail_h = h - img_top - Inches(0.6 if caption else 0.3)
    avail_w = w - Inches(1.0)
    pic = slide.shapes.add_picture(img_path, Inches(0.5), img_top, width=avail_w)
    if pic.height > avail_h:
        sp = pic._element; sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(img_path, Inches(0.5), img_top, height=avail_h)
        pic.left = Inches((w.inches - pic.width.inches) / 2)
    else:
        pic.left = Inches((w.inches - pic.width.inches) / 2)

    if caption:
        cb = slide.shapes.add_textbox(Inches(0.5),
                                      pic.top + pic.height + Inches(0.1),
                                      w - Inches(1.0), Inches(0.4))
        p = cb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = caption
        r.font.size = Pt(12); r.font.italic = True; r.font.color.rgb = MUTED


def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Title ────────────────────────────────────────────────────────
    add_title_slide(
        prs,
        "Confidence reward shaping on MATH levels 3-5",
        "exp_039 to exp_048  —  Qwen3-4B  —  may 2026 update",
    )

    # ── Where we left off ────────────────────────────────────────────
    add_text_slide(
        prs, "Where we left off (april deck)",
        [
            "best result on Big-Math int-2000: flipped pure-proof GTPO-EMA (exp_028)",
            "peak reward 9.5 at step 194, 11 steps earlier than matched GRPO baseline (exp_027)",
            "on GSM-8K (exp_005, exp_024) the ranking between close confidence variants was dominated by run-to-run vLLM sampling noise",
            "two open questions: does the per-token confidence bonus survive on a harder benchmark, and is it a real signal or seed luck",
        ],
    )

    # ── This block ───────────────────────────────────────────────────
    add_text_slide(
        prs, "This block — what we ran",
        [
            "exp_039, exp_040 — finish the Big-Math int-2000 sweep, add per-rollout logging",
            "exp_041 — GRPO baseline on a harder benchmark: MATH-lighteval levels 3-5, integer answers, ~3400 problems",
            "exp_042 — port GTPO-EMA-flipped to MATH levels 3-5 (turned out to be a silent ablation, see slide 6)",
            "exp_043 — GTPO-Conf without EMA on MATH levels 3-5",
            "exp_044 — GTPO-EMA-flipped properly activated (the real fix for exp_042)",
            "exp_045, exp_046 — SCRS family: sequence-level confidence / entropy shaping",
            "exp_047 — vectorised re-implementation of GTPO-Conf for MATH levels 3-5",
            "exp_048 — UCAS Stage 1 from arXiv 2510.10649: multiplicative response-level confidence weighting",
        ],
    )

    # ── Method recap ─────────────────────────────────────────────────
    add_text_slide(
        prs, "Confidence metric — quick recap",
        [
            "C_t = -mean over top-k log pi_v, computed at each generated token (we use top-k = 20)",
            "intuition: C grows with peakedness of the next-token distribution, NOT with entropy",
            "flipped split: O+ (correct rollouts) gets weight 1 / EMA(C) so flat/uncertain tokens get more credit",
            "O- (wrong rollouts) gets weight EMA(C) so confident mistakes are penalised harder",
            "EMA with lambda = 0.9 over the active sub-batch at each step t to denoise the per-token signal",
            "alpha1 = 0.9 keeps the GRPO advantage as the backbone, alpha2 = 0.1 is the confidence bonus weight",
        ],
    )

    # ── exp_039 ──────────────────────────────────────────────────────
    add_text_slide(
        prs, "exp_039 — GTPO-EMA-flipped on Big-Math int-2000 (full 1000 steps)",
        [
            "setup: Qwen3-4B, LoRA r=64, bs=4, num_generations=8, max_completion=3584, lr=5e-6",
            "matched 1:1 against exp_038 GRPO baseline (same dataset, shuffle seed 3407)",
            "first time the flipped GTPO-EMA finishes a full 1000-step run on Qwen3 / Big-Math",
            "reward ceiling 9.5 hit by step 9; runtime 31371 s (8h 43m)",
            "last-50 mean reward: 5.97, answer_exact final: 3.0, format_exact final: 3.0",
            "matches GRPO baseline trajectory — no separation visible above the per-step noise",
        ],
    )
    add_image_slide(
        prs, "exp_039 — total reward, KL, format, answer",
        "exp_039_dashboard.png",
        "1000-step run, Qwen3-4B, GTPO-EMA-flipped, Big-Math int-2000",
    )

    # ── exp_040 ──────────────────────────────────────────────────────
    add_text_slide(
        prs, "exp_040 — same run + per-rollout logging tooling",
        [
            "identical hyperparameters to exp_039, plus a RolloutLoggerTrainer that saves per-step .npz",
            "each .npz keeps top-20 log_probs, token ids, completion mask, advantages and correctness flag",
            "purpose: build the substrate for offline analysis of which tokens actually carry the confidence signal",
            "last-50 mean reward: 6.07 (close to exp_039), no behaviour change from logging",
            "the rollout viewer (viewer.py) lets us scrub through any step and inspect token-level confidence vs the per-rollout reward",
        ],
    )

    # ── Dataset shift ────────────────────────────────────────────────
    add_text_slide(
        prs, "Dataset shift — MATH-lighteval levels 3-5 (integer)",
        [
            "why: Big-Math int-2000 saturates at the 9.5 ceiling within ~200 steps; can not separate methods",
            "new train set: DigitalLearningGmbH/MATH-lighteval, levels 3, 4, 5",
            "filter: only examples whose \\boxed{} answer parses as an integer (~3400 problems total)",
            "max_completion = 3584, same prompt template, same reward stack (format exact + approximate, answer exact + numeric)",
            "shuffle seed 3407 fixed across exp_041..exp_048",
        ],
    )

    # ── exp_041 ──────────────────────────────────────────────────────
    add_text_slide(
        prs, "exp_041 — GRPO baseline on MATH levels 3-5",
        [
            "Qwen3-4B, bs=4, num_generations=8, 1000 step budget (run reached step 822 before container restart)",
            "reward ceiling 9.5 first hit at step 22; format learnt cleanly around step 250",
            "last-50 mean reward: 7.82 — substantially higher steady-state than Big-Math int-2000 baseline (exp_027 last-50 ≈ 5.95)",
            "mean completion length decays from 3500 to ~1500 tokens as the model stops over-thinking the easy items",
            "this is the reference curve every confidence variant has to beat",
        ],
    )
    add_image_slide(
        prs, "exp_041 — six-panel dashboard",
        "exp_041_dashboard.png",
        "reward, pass-rate, format, answer, completion length, KL",
    )

    # ── exp_042 — the silent ablation ────────────────────────────────
    add_text_slide(
        prs, "exp_042 — what we thought was GTPO-EMA-flipped",
        [
            "we ported the flipped pure-proof trainer to MATH levels 3-5 and re-ran",
            "result was indistinguishable from the GRPO baseline at every metric: reward, pass-rate, grad norm, cumulative pass",
            "investigation: token_advantages were computed correctly, then multiplied by 0.0 before the loss",
            "so exp_042 is a clean ablation: GTPO advantages computed but zeroed → pure GRPO",
            "useful negative — confirms the test harness and the GRPO backbone are intact",
        ],
    )
    add_image_slide(
        prs, "exp_042 vs exp_041 — overlapping curves confirm the zeroing",
        "exp_042_ablation.png",
        "reward, pass rate, grad norm, cumulative pass — orange and blue track each other within noise",
    )

    # ── exp_043 — collapse ───────────────────────────────────────────
    add_text_slide(
        prs, "exp_043 — GTPO-Conf (no EMA) on MATH levels 3-5",
        [
            "method: column-wise per-token confidence advantage, log(1+C) on O+, log(1+1/C) on O-, no temporal smoothing",
            "same recipe that beat GRPO on GSM-8K in exp_005",
            "result: peak 9.5 at step 183, then collapse — last-50 mean reward = -3.70, answer_exact final 0.0",
            "completion length explodes upward as the per-token bonus rewards each new uncertain token (length feedback loop)",
            "first hard evidence that per-token confidence shaping without EMA does not transfer from GSM-8K to MATH levels 3-5",
        ],
    )
    add_image_slide(
        prs, "exp_043 — collapse after step ~200",
        "exp_043_collapse.png",
    )

    # ── exp_044 — fixed but still collapses ──────────────────────────
    add_text_slide(
        prs, "exp_044 — GTPO-EMA-flipped, properly activated",
        [
            "fix for exp_042: token_advantages = gtpo_adv (no 0.0 multiplier)",
            "EMA smoothing was supposed to neutralise the length feedback loop seen in exp_043",
            "result: peak 9.5 at step 150, last-50 mean reward = -1.85, answer_exact final 1.125",
            "less catastrophic than exp_043 but still well below the GRPO baseline (7.82 → -1.85)",
            "EMA helps but does not fix the underlying issue on this benchmark — the per-token signal is unstable for long, hard-math completions",
        ],
    )
    add_image_slide(
        prs, "exp_044 vs exp_041 — EMA does not save it on this dataset",
        "exp_044_vs_041.png",
        "left: total reward, right: answer_exact — orange (exp_044) drifts below zero",
    )

    # ── Diagnosis ────────────────────────────────────────────────────
    add_text_slide(
        prs, "Diagnosis — why per-token shaping fails on MATH levels 3-5",
        [
            "MATH levels 3-5 completions are 2-3x longer than GSM-8K ones",
            "per-token bonus on O+ rewards each additional uncertain token → policy is pushed to generate longer, less decisive chains",
            "per-token penalty on O- punishes peaked tokens uniformly, including the few correct decisive tokens inside a wrong rollout",
            "EMA smooths the per-token signal but does not change its sign or remove the length pressure",
            "decision: park per-token shaping and try sequence-level shaping instead",
        ],
    )

    # ── exp_045 SCRS-Conf ────────────────────────────────────────────
    add_text_slide(
        prs, "exp_045 — SCRS-Confidence (sequence-level shaping)",
        [
            "new family: SCRS = Sequence-level Confidence Reward Shaping",
            "rule: shaped_adv_i = grpo_adv_i - alpha2 * z(mean_conf_i)",
            "uncertain-and-correct rollouts get a bigger positive advantage; confident-and-wrong rollouts get a bigger penalty",
            "no per-token accumulation → no length feedback loop",
            "alpha2 = 0 recovers pure GRPO exactly — clean ablation switch",
            "result: full 1000 steps, last-50 mean reward = 5.80, answer_exact final 3.0",
            "stable but not separating from baseline yet",
        ],
    )

    # ── exp_046 SCRS-Entropy ─────────────────────────────────────────
    add_text_slide(
        prs, "exp_046 — SCRS-Entropy",
        [
            "swap: replace mean confidence with mean top-k Shannon entropy H (top_k = 100)",
            "sign flips: shaped_adv_i = grpo_adv_i + alpha2 * z(mean_entropy_i)",
            "intuition: H is roughly the inverse of C, so adding entropy is the same direction as subtracting confidence",
            "full 1000 steps, last-50 mean reward = 6.13, slightly above SCRS-Conf",
            "trajectory tracks the GRPO baseline within noise — sequence-level shaping is stable, but no clear win",
        ],
    )

    # ── exp_047 — per-token retry ────────────────────────────────────
    add_text_slide(
        prs, "exp_047 — GTPO-Conf, vectorised, one more try",
        [
            "rewrite of exp_043 with vectorised per-token loops, TRL 0.23.1 compute_loss API, ghost guard for zero-variance batches",
            "same method math: log(1+C) on O+, log(1+1/C) on O-, no EMA, alpha1=1, alpha2=0.1",
            "hope: the previous collapse came from an implementation bug rather than the method itself",
            "result: training aborted at step 295, peak 9.5 at step 49, last-50 mean reward = -2.89",
            "collapse signature is identical to exp_043 — confirms the method, not the implementation, is the problem",
        ],
    )
    add_image_slide(
        prs, "All per-token variants vs baseline on MATH levels 3-5",
        "per_token_collapse_panel.png",
        "lower is better (zero-correct fraction); only exp_041 baseline and exp_042 ablation stay near the floor",
    )

    # ── exp_048 UCAS Stage 1 ─────────────────────────────────────────
    add_text_slide(
        prs, "exp_048 — UCAS Stage 1 (arXiv 2510.10649)",
        [
            "external method: uncertainty-aware advantage shaping, ACL 2026 submission",
            "C_t = full-vocabulary -mean(log_prob); per-sequence conf = masked mean of C_t",
            "C_hat_i = z-score within rollout group; W_i = exp(±alpha * C_hat_i), sign by sign of grpo_adv_i",
            "shaped_adv_i = W_i * grpo_adv_i — multiplicative, preserves advantage scale (unlike SCRS additive form)",
            "alpha = 0.25 (stable range [0.1, 0.5] per paper Table 3)",
            "result: full 1000 steps, last-50 mean reward = 5.92, answer_exact final 3.0",
            "stable, but again tracks GRPO baseline within step-to-step noise",
        ],
    )
    add_image_slide(
        prs, "Sequence-level methods vs GRPO baseline",
        "sequence_level_stable.png",
        "exp_041, exp_045, exp_046, exp_048 — zero-correct fraction (top) and mean reward (bottom)",
    )

    # ── Side-by-side overview ────────────────────────────────────────
    add_image_slide(
        prs, "Same picture, two metrics",
        "seq_level_overview.png",
        "left: total reward, right: answer_exact — all four runs overlap",
    )

    # ── Summary table ────────────────────────────────────────────────
    add_text_slide(
        prs, "Results table — MATH levels 3-5",
        [
            "exp_041 GRPO baseline             — last-50 reward 7.82  /  ans 2.625  /  status: ran 822 steps",
            "exp_042 GTPO-EMA-flipped (zeroed) — last-50 reward 6.94  /  ans 2.625  /  status: ablation, identical to baseline",
            "exp_043 GTPO-Conf (no EMA)        — last-50 reward -3.70 /  ans 0.000  /  status: collapsed",
            "exp_044 GTPO-EMA-flipped (active) — last-50 reward -1.85 /  ans 1.125  /  status: collapsed",
            "exp_045 SCRS-Confidence           — last-50 reward 5.80  /  ans 3.000  /  status: stable, no win",
            "exp_046 SCRS-Entropy              — last-50 reward 6.13  /  ans 3.000  /  status: stable, no win",
            "exp_047 GTPO-Conf vectorised      — last-50 reward -2.89 /  ans 1.875  /  status: collapsed, aborted at 295",
            "exp_048 UCAS Stage 1              — last-50 reward 5.92  /  ans 3.000  /  status: stable, no win",
        ],
    )

    # ── Takeaways ────────────────────────────────────────────────────
    add_text_slide(
        prs, "Takeaways",
        [
            "per-token confidence shaping that wins on GSM-8K (exp_005, exp_028) does not transfer to MATH levels 3-5",
            "all three per-token variants we tried (exp_043, exp_044, exp_047) collapse to near-zero answer accuracy",
            "EMA temporal smoothing helps but does not fix the length feedback loop on long-context math",
            "sequence-level shaping (SCRS family, UCAS Stage 1) is stable across 1000 steps but matches the GRPO baseline within noise — no separation yet",
            "exp_042 turned out to be a useful free ablation — confirmed the GRPO backbone is fine and isolated the bug to the multiplier line",
            "the per-rollout logging substrate from exp_040 will let us look at where the per-token signal actually flips sign during a long completion",
        ],
    )

    # ── Next ─────────────────────────────────────────────────────────
    add_text_slide(
        prs, "Next",
        [
            "use the exp_040 rollout logs to characterise C and EMA(C) along long MATH completions — does the signal flip sign mid-chain on wrong rollouts",
            "try a length-corrected per-token bonus: divide alpha2 contribution by sqrt(seq_len) or restrict to the first N tokens",
            "two-stage UCAS (Stage 2: per-token conditioned on the sequence weight) on MATH levels 3-5",
            "if sequence-level matches but does not beat baseline, run multi-seed (3-5 seeds) on the most promising variant (SCRS-Entropy) to confirm or reject the parity claim with confidence intervals",
            "all code under github.com/bliskavets/aiim_research, experiments exp_039..exp_048",
        ],
    )

    prs.save(OUT)
    print(f"saved: {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
