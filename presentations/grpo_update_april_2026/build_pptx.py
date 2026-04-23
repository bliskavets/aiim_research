"""
Build grpo_update_april_2026.pptx.

Opens GRPO_Modifications_Report.pptx, keeps slides 1..11 (title + Overview,
Methods, Progression, exp_001/002/004/005 description+results), then
appends new slides covering exp_021..028 with exp_026 first (proof-based
flipped confidence that works).

Run:   python3 build_pptx.py
Output: grpo_update_april_2026.pptx
"""
import copy, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
SRC  = os.path.join(HERE, "..", "..", "experiments", "GRPO_Modifications_Report.pptx")
IMG  = os.path.join(HERE, "img")
OUT  = os.path.join(HERE, "grpo_update_april_2026.pptx")

# ── Palette (same as the source deck) ─────────────────────────────────
ACCENT   = RGBColor(0x4f, 0x46, 0xe5)
ACCENT2  = RGBColor(0x02, 0x84, 0xc7)
TEXT     = RGBColor(0x1e, 0x29, 0x3b)
MUTED    = RGBColor(0x64, 0x74, 0x8b)
GREEN    = RGBColor(0x05, 0x96, 0x69)
ORANGE   = RGBColor(0xd9, 0x77, 0x06)


# ──────────────────────────────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────────────────────────────

def keep_only_first_n_slides(prs: Presentation, n: int) -> None:
    """Mutate prs so only the first n slides remain."""
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    # Remove from the end to preserve indices of the kept ones
    for sid in ids[n:]:
        rId = sid.get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sid)


def add_title_and_body(
    prs: Presentation,
    title: str,
    body_lines: list[str],
    img_path: str | None = None,
    caption: str | None = None,
) -> None:
    """
    Append a slide with a title + bullet body + optional image below.
    Uses the blank layout so we control the whole surface.
    """
    blank_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank_layout)

    # Try to strip the layout's default placeholders if they show up
    for shape in list(slide.shapes):
        if shape.has_text_frame and shape.text_frame.text.strip() in ("Click to edit Master title style", "Title"):
            sp = shape._element
            sp.getparent().remove(sp)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Title strip
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3),
                                         slide_w - Inches(1.0), Inches(0.7))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = ACCENT

    # Body
    body_top = Inches(1.15)
    body_h = Inches(2.0) if img_path else Inches(5.8)
    body_box = slide.shapes.add_textbox(Inches(0.5), body_top,
                                        slide_w - Inches(1.0), body_h)
    tf = body_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(16)
        run.font.color.rgb = TEXT

    # Optional image
    if img_path and os.path.exists(img_path):
        # Fit below body, keep aspect by setting only one dim
        img_top = body_top + body_h + Inches(0.1)
        avail_h = slide_h - img_top - Inches(0.2)
        avail_w = slide_w - Inches(1.0)
        pic = slide.shapes.add_picture(
            img_path,
            Inches(0.5), img_top,
            width=avail_w,
        )
        # If picture overflows vertical space, fall back to height-fit
        if pic.height > avail_h:
            # Remove oversized one, re-add with height-fit
            sp = pic._element
            sp.getparent().remove(sp)
            pic = slide.shapes.add_picture(
                img_path,
                Inches(0.5), img_top,
                height=avail_h,
            )
            # Center horizontally
            pic.left = Inches((prs.slide_width.inches - pic.width.inches) / 2)

        if caption:
            cap = slide.shapes.add_textbox(
                pic.left, pic.top + pic.height + Inches(0.05),
                pic.width, Inches(0.3),
            )
            p = cap.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = caption
            run.font.size = Pt(11)
            run.font.italic = True
            run.font.color.rgb = MUTED


def add_image_slide(prs: Presentation, title: str, img_path: str,
                    caption: str | None = None) -> None:
    """Slide with just a title and one large centered image, no stretch."""
    blank_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank_layout)
    for shape in list(slide.shapes):
        if shape.has_text_frame and shape.text_frame.text.strip() in ("Click to edit Master title style", "Title"):
            sp = shape._element
            sp.getparent().remove(sp)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3),
                                         slide_w - Inches(1.0), Inches(0.7))
    p = title_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = ACCENT

    # Image area
    img_top = Inches(1.15)
    avail_h = slide_h - img_top - Inches(0.6 if caption else 0.3)
    avail_w = slide_w - Inches(1.0)

    pic = slide.shapes.add_picture(img_path, Inches(0.5), img_top, width=avail_w)
    if pic.height > avail_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(img_path, Inches(0.5), img_top, height=avail_h)
        pic.left = Inches((slide_w.inches - pic.width.inches) / 2)
    else:
        pic.left = Inches((slide_w.inches - pic.width.inches) / 2)

    if caption:
        cap = slide.shapes.add_textbox(Inches(0.5), pic.top + pic.height + Inches(0.1),
                                       slide_w - Inches(1.0), Inches(0.4))
        p = cap.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = caption
        r.font.size = Pt(12)
        r.font.italic = True
        r.font.color.rgb = MUTED


# ──────────────────────────────────────────────────────────────────────
# Build
# ──────────────────────────────────────────────────────────────────────

def main() -> None:
    prs = Presentation(SRC)
    print(f"source: {len(prs.slides)} slides")

    # Keep slides 1..11 (index 0..10)
    keep_only_first_n_slides(prs, 11)
    print(f"after trim: {len(prs.slides)} slides")

    # ── HEADLINE: exp_026 proof-based flipped works ──────────────────
    add_title_and_body(
        prs,
        "exp 026: proof based GRPO + flipped confidence on GSM-8K",
        [
            "setup: GSM8K, Llama-3.2-3B-Instruct, LoRA r=64",
            "bs=1, grad_accum=4, num_generations=4, 500 steps, lr=5e-6",
            "method: pure-proof GTPO-EMA from experiments/proof/GTPO-EMA-full.txt",
            "key idea: we noticed that C = mean log pi over top-k actually grows with peakedness, not with entropy",
            "so we swap the signal roles between O+ and O-. O+ now rewards flat tokens (exploration on correct paths), O- penalizes peaked tokens (confident mistakes)",
            "result: peak reward 9.5 at step 358, final reward 3.0, format_exact 3.0",
            "KL at step 500: 0.095, the lowest among successful confidence methods (exp_005 was 0.069 but GRPO-S-Conf collapsed there)",
        ],
    )
    add_image_slide(
        prs,
        "exp 026 vs GRPO baseline on GSM-8K",
        os.path.join(IMG, "exp_026_flipped_vs_grpo.png"),
        caption="flipped pure-proof GTPO-EMA reaches the 9.5 reward ceiling with low KL",
    )

    # ── exp_021: GTPO-Conf on Big-Math ────────────────────────────────
    add_title_and_body(
        prs,
        "exp 021: GTPO + confidence bonus (no EMA) on Big-Math integer",
        [
            "setup: Big-Math-RL-Verified integer subset, Llama-3.2-3B",
            "bs=4, num_generations=16, 1000 steps, lr=5e-6",
            "max_completion_length=3072, bf16",
            "method: confidence based bonus, no EMA smoothing",
            "result: peak reward 9.5 at step 230, stable format by step 200",
            "confidence without EMA is enough on Big-Math as well, not just GSM8K",
        ],
    )
    add_image_slide(
        prs,
        "exp 021 vs exp 017 baseline on Big-Math",
        os.path.join(IMG, "exp_021_compare.png"),
    )

    # ── exp_022: GTPO binary ──────────────────────────────────────────
    add_title_and_body(
        prs,
        "exp 022: GTPO with binary O+ / O- split",
        [
            "setup: same as exp 021 on Big-Math integer",
            "method: binary split by reward_answer_exact >= 0",
            "O+ contains all non-wrong samples including no-format and within-20 percent",
            "O- contains only wrong-answer-in-format (-1.5)",
            "result: peak reward 9.5 at step 202",
            "cleaner signal than z-score advantages, but threshold was too permissive",
        ],
    )
    add_image_slide(
        prs,
        "exp 022 vs exp 017 baseline on Big-Math",
        os.path.join(IMG, "exp_022_compare.png"),
    )

    # ── exp_023: GTPO-EMA binary ──────────────────────────────────────
    add_title_and_body(
        prs,
        "exp 023: GTPO-EMA with binary O+ / O- split",
        [
            "setup: same as exp 022 on Big-Math integer",
            "method: EMA smoothed confidence + binary O+ / O- split",
            "result: peak reward 9.5 at step 215",
            "combines benefits of EMA smoothing with cleaner binary split",
        ],
    )
    add_image_slide(
        prs,
        "exp 023 vs exp 017 baseline on Big-Math",
        os.path.join(IMG, "exp_023_compare.png"),
    )

    # ── exp_024: reproduction of exp_005 ──────────────────────────────
    add_title_and_body(
        prs,
        "exp 024: byte-identical replay of exp 005 on GSM-8K",
        [
            "motivation: check if the exp 005 win over GRPO baseline was real or lucky",
            "setup: GSM8K, Llama-3.2-3B, identical code and seed as exp 005",
            "bs=1, grad_accum=4, num_generations=4, 500 steps",
            "result: the ranking flipped between runs",
            "exp 005: GTPO-Conf won (9.5), GRPO-S-Conf collapsed (2.0)",
            "exp 024: GTPO-Conf collapsed (3.1), GRPO-S-Conf won (9.5)",
            "conclusion: run-to-run variance from vLLM sampling is larger than the method gap at this scale",
        ],
    )
    add_image_slide(
        prs,
        "exp 024 GTPO-Conf reproduction on GSM-8K",
        os.path.join(IMG, "exp_024_gtpo_conf.png"),
        caption="same code as exp 005 but collapse happens instead of reaching the ceiling",
    )
    add_image_slide(
        prs,
        "exp 024 GRPO-S-Conf reproduction on GSM-8K",
        os.path.join(IMG, "exp_024_grpos_conf.png"),
        caption="same code as exp 005 but ceiling reached instead of collapse",
    )

    # ── exp_025: pure-proof on GSM-8K ─────────────────────────────────
    add_title_and_body(
        prs,
        "exp 025: pure-proof GTPO-EMA on GSM-8K",
        [
            "setup: GSM8K, Llama-3.2-3B, same hyperparameters as exp 005",
            "method: implement Def 1.4 of GTPO-EMA-full.txt literally",
            "per-token bonus weighted by EMA(C) over active sequences in O+ at each step t, multiplied by d_t",
            "conservation alpha1+alpha2=1 so total reward mass stays constant",
            "separate z-norm on O+ and O- as in Def 1.5",
            "result: peak 9.5 at step 253, reward 3.0, format_exact 3.0",
            "matches the successful cluster but not yet distinguishable from run-to-run variance",
        ],
    )
    add_image_slide(
        prs,
        "exp 025 vs GRPO baseline on GSM-8K",
        os.path.join(IMG, "exp_025_pureproof.png"),
    )

    # ── exp_027: new Big-Math int 2000 baseline ───────────────────────
    add_title_and_body(
        prs,
        "exp 027: GRPO baseline on Big-Math integer 2000",
        [
            "setup: Big-Math integer filter, 2000 shuffled examples (seed 3407)",
            "Llama-3.2-3B, bs=4, grad_accum=1, num_generations=8",
            "500 steps, max_completion_length=2048, lr=5e-6",
            "result: peak reward 9.5 at step 205",
            "last 50 steps average: reward 5.95, format_exact 2.79",
            "clear S-curve, format learned near step 170, stable high band 4 to 8 after",
        ],
    )
    add_image_slide(
        prs,
        "exp 027 reward progress on Big-Math integer 2000",
        os.path.join(IMG, "exp_027_progress.png"),
    )

    # ── exp_028: flipped pure-proof on Big-Math int 2000 ─────────────
    add_title_and_body(
        prs,
        "exp 028: flipped pure-proof GTPO-EMA on Big-Math integer 2000",
        [
            "setup: identical to exp 027 (same 2000 shuffled integer problems, same hyperparameters)",
            "method: flipped pure-proof from exp 026",
            "O+ / O- split by reward_answer_exact >= 1.0 via a reward_cache module",
            "O+ covers exact, strip and within-10 percent matches",
            "O- covers within-20 percent, no-format and wrong answer",
            "result so far: peak 9.5 at step 194, 11 steps earlier than baseline (205)",
            "stable format_exact=3.0 and higher answer_exact than baseline in the running snapshot",
        ],
    )
    add_image_slide(
        prs,
        "exp 028 vs exp 027 on Big-Math integer 2000",
        os.path.join(IMG, "exp_028_vs_027.png"),
        caption="green is exp 028 flipped pure-proof, gray is exp 027 GRPO baseline",
    )

    # ── Summary ───────────────────────────────────────────────────────
    add_title_and_body(
        prs,
        "Summary of new work (exp 021 to 028)",
        [
            "exp 026 is the first confidence-based GRPO variant we built directly from the proof, and it works on GSM-8K",
            "the swap of EMA vs 1/EMA between O+ and O- is motivated by a numerical fact about the top-k confidence metric",
            "exp 024 showed that ranking between close confidence-variants is dominated by run-to-run variance on GSM-8K",
            "exp 021 to 023 confirmed the confidence bonus and binary O+ / O- split scale from GSM-8K to Big-Math integer",
            "exp 027 and 028 move to Big-Math integer 2000 as a harder benchmark",
            "early result: exp 028 flipped reaches the reward ceiling 11 steps before the matched GRPO baseline",
            "all code and figures live in github.com/bliskavets/aiim_research",
        ],
    )

    prs.save(OUT)
    print(f"saved: {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
