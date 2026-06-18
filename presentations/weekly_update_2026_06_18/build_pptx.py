"""
build_pptx.py — weekly update deck (2026-06-11 .. 2026-06-18).
One slide per experiment run this week: short title (<=7 words), the main figure,
a one-line takeaway. Light / minimal style. Covers BOTH this session's
experiments and the parallel Claude session's (exp_056, exp_058, exp_059-instruct).

Run: python build_pptx.py  -> weekly_update_2026_06_18.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
EXP = os.path.normpath(os.path.join(HERE, "..", "..", "experiments"))

# palette — light & clean
INK   = RGBColor(0x1f, 0x29, 0x37)   # near-black slate
GRAY  = RGBColor(0x6b, 0x72, 0x80)   # muted gray
ACCENT = RGBColor(0x4f, 0x46, 0xe5)  # indigo
WHITE = RGBColor(0xff, 0xff, 0xff)
RULE  = RGBColor(0xe2, 0xe8, 0xf0)

SW, SH = Inches(13.333), Inches(7.5)   # 16:9

# (kicker, title<=7w, image rel-path or None, takeaway)
SLIDES = [
    ("exp_056 · Search-R1 (NQ+HotpotQA, wiki-18)",
     "Search-R1: shaping never beats GRPO",
     "exp_056_searchr1_qwen3_grpo_vs_shaped/figures/exp056_fixed_4way_final.png",
     "All 4 methods land at EM ≈ 0.196; the 3 shaped curves are ~identical. "
     "gtpo_conf KL explodes (~1e7, grad ~5e8) and grpo_s blows up — shaped-advantage magnitudes uncontrolled."),

    ("exp_057 · Qwen3-4B instruct · Omni-MATH integer",
     "Shaping was silently bypassed — fixed",
     "exp_057_qwen3_native_omnimath_int_4way/figures/exp057_progress.png",
     "Audit found unsloth's compiled loss bypassed the shaping (shaped runs = plain GRPO). "
     "After the injection fix, per-token shaping (gtpo_conf) DEGRADES vs GRPO (L50 −0.2 vs +2.6)."),

    ("exp_058 · Qwen3-4B-BASE · Big-Math int-2000  [parallel]",
     "Base model: gtpo_ema collapses, others learn",
     "exp_058_qwen3base_bigmath_int2k_4way/figures/exp058_4way_base_model.png",
     "Base has headroom: grpo & grpo_s_entropy learn (boxed 0.7→1.9); gtpo_conf lags; "
     "gtpo_ema_flipped COLLAPSES via length-explosion (640 → 3400 tok, boxed → 0)."),

    ("exp_059 · Qwen3-4B-BASE · GSM8K (exp_005 candidates)",
     "Per-token confidence hurts; seq-level ties",
     "exp_059_exp005candidates_qwen3base_gsm8k/figures/exp059_progress.png",
     "grpo +8.66, grpo_s_conf +8.84 (tie, within noise), gtpo_conf +0.86 (never learns the format, KL≈0.002). "
     "Inverts exp_005's old GTPO-Conf>GRPO-S claim (that run was bypassed)."),

    ("exp_059 · Qwen3-4B-instruct · hard Big-Math  [parallel]",
     "Hard Big-Math 4-way — scaffolded, pending",
     None,
     "Parallel session: 4 FIXED-framework methods on the hard Big-Math slice (llama8b<0.3), Qwen3-4B-instruct. "
     "Code + tests committed; training run still pending (no results yet)."),

    ("exp_060 · CONTROL · GRPO-S with β2 = 0",
     "Control: β2=0 reproduces GRPO (code valid)",
     "exp_060_grpos_beta2zero_control/figures/exp060_control_stepmatched.png",
     "With the entropy bonus off, GRPO-S tracks GRPO step-for-step (reward +0.67 vs +0.64, grad matches). "
     "Confirms the GRPO-S injection / gradient path is correct — the harm is the shaping itself, not a bug."),
]

SUMMARY = [
    "Found & fixed a real bug: unsloth's compiled GRPO loss silently bypassed every shaping trainer's "
    "_compute_loss — \"shaped\" runs were plain GRPO. Likely affects the earlier exp_049→056 arc.",
    "With shaping ACTUALLY applied, no per-token method beats GRPO — across Qwen3 instruct & base, "
    "GSM8K, Big-Math and Search-R1 it ties at best and usually hurts (gtpo_ema even collapses).",
    "Sequence-level shaping (grpo_s) is benign (ties GRPO); per-token shaping is the harmful one "
    "(reward-misaligned z-normed advantage; uncontrolled magnitudes / KL explosions).",
    "Control (exp_060) proves the code is correct — the historical \"shaping helps\" results were "
    "almost certainly measured with the shaping bypassed.",
]


def _txt(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = font
    return tb


def _rule(slide, l, t, w):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(3))
    sh.fill.solid(); sh.fill.fore_color.rgb = ACCENT; sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def add_blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    return s


def title_slide(prs):
    s = add_blank(prs)
    _txt(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.2),
         "Reward-Shaping in GRPO — Weekly Review", 40, INK, bold=True)
    _rule(s, Inches(0.95), Inches(3.55), Inches(3.2))
    _txt(s, Inches(0.95), Inches(3.8), Inches(11), Inches(0.8),
         "Qwen3-4B · GSM8K / Big-Math / Omni-MATH / Search-R1 · 11–18 June 2026", 18, GRAY)
    _txt(s, Inches(0.95), Inches(6.7), Inches(11.5), Inches(0.5),
         "6 experiments across two parallel sessions", 13, GRAY)
    return s


def exp_slide(prs, kicker, title, img, takeaway):
    s = add_blank(prs)
    _txt(s, Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.35), kicker, 12, ACCENT, bold=True)
    _txt(s, Inches(0.6), Inches(0.66), Inches(12.1), Inches(0.8), title, 28, INK, bold=True)
    _rule(s, Inches(0.62), Inches(1.5), Inches(2.6))
    # figure centered in the body
    if img:
        path = os.path.join(EXP, img)
        if os.path.exists(path):
            from PIL import Image
            iw, ih = Image.open(path).size
            maxw, maxh = Inches(11.6), Inches(4.5)
            scale = min(maxw / iw, maxh / ih)
            w, h = Emu(int(iw * scale)), Emu(int(ih * scale))
            left = Emu(int((SW - w) / 2)); top = Inches(1.75)
            s.shapes.add_picture(path, left, top, width=w, height=h)
        else:
            _txt(s, Inches(0.8), Inches(3.0), Inches(11.6), Inches(1.0),
                 f"[figure missing: {img}]", 14, GRAY, align=PP_ALIGN.CENTER)
    else:
        _txt(s, Inches(0.8), Inches(3.0), Inches(11.6), Inches(1.2),
             "code + tests committed · training run pending", 20, GRAY, align=PP_ALIGN.CENTER)
    # takeaway band at the bottom
    _txt(s, Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.95), takeaway, 14, INK)
    return s


def summary_slide(prs):
    s = add_blank(prs)
    _txt(s, Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.8),
         "Takeaway: applied shaping doesn't beat GRPO", 30, INK, bold=True)
    _rule(s, Inches(0.62), Inches(1.35), Inches(2.6))
    tb = s.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = tb.text_frame; tf.word_wrap = True
    for i, b in enumerate(SUMMARY):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        r = p.add_run(); r.text = "•  " + b
        f = r.font; f.size = Pt(17); f.color.rgb = INK; f.name = "Calibri"
    return s


def formula_slide(prs, title, blocks):
    """blocks: list of (label_or_None, [lines]). Minimal words, mono formulas."""
    s = add_blank(prs)
    _txt(s, Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.8), title, 28, INK, bold=True)
    _rule(s, Inches(0.62), Inches(1.35), Inches(2.6))
    tb = s.shapes.add_textbox(Inches(0.8), Inches(1.75), Inches(11.8), Inches(5.4))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for label, lines in blocks:
        if label is not None:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.space_before = Pt(0 if first else 12); p.space_after = Pt(3)
            r = p.add_run(); r.text = label
            f = r.font; f.size = Pt(15); f.bold = True; f.color.rgb = ACCENT; f.name = "Calibri"
            first = False
        for ln in lines:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.space_after = Pt(3)
            r = p.add_run(); r.text = ln
            f = r.font; f.size = Pt(16); f.color.rgb = INK; f.name = "Consolas"
            first = False
    return s


def main():
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH
    title_slide(prs)
    for kicker, title, img, takeaway in SLIDES:
        exp_slide(prs, kicker, title, img, takeaway)
    summary_slide(prs)

    # --- method-difference slide ---
    formula_slide(prs, "gtpo_conf  vs  gtpo_ema_flipped", [
        (None, ["both:   Cₜ = −mean₍top-k₎ log π(v | ·)        # token confidence",
                "both:   separate z-norm over O⁺ / O⁻ tokens"]),
        ("gtpo_conf  — raw, monotone", [
            "O⁺:  r̃ₜ = α₁ + α₂·(C̃ₜ / Σ C̃)·dₜ ,   C̃ = log(1+C)",
            "O⁻:  r̃ₜ = −[α₁ + α₂·(Ĩₜ / Σ Ĩ)·hₜ] ,  Ĩ = log(1+1/C)"]),
        ("gtpo_ema_flipped  — EMA-smoothed + roles swapped", [
            "EMAₜ = λ·EMAₜ₋₁ + (1−λ)·Cₜ",
            "O⁺ weight = 1/EMAₜ   (reward LOW-confidence tokens)",
            "O⁻ weight =  EMAₜ     (penalize HIGH-confidence tokens)"]),
        ("difference", [
            "conf = instantaneous C, high-C → high bonus",
            "ema  = smoothed C  +  inverted O⁺/O⁻ weighting"]),
    ])

    # --- length-penalty hack slide ---
    formula_slide(prs, "Length-penalty hack (gtpo_ema_flipped)", [
        ("why it explodes (exp_058)", [
            "O⁺ bonus ∝ 1/EMA(C) → rewards low-confidence tokens",
            "→ model farms length:  640 → 3400 tok,  boxed → 0"]),
        ("fix 1 — soft length penalty on seq reward (before O⁺/O⁻ split)", [
            "rᵢ ← rᵢ − γ · max(0, |oᵢ| − L₀) / L₀        (γ≈0.5, L₀ = target len)"]),
        ("fix 2 — length-normalize the bonus (can't grow with length)", [
            "bonusₜ ← bonusₜ · L₀ / max(L₀, |oᵢ|)"]),
        ("effect", [
            "caps exploration-reward per completion → removes the incentive to ramble"]),
    ])
    out = os.path.join(HERE, "weekly_update_2026_06_18.pptx")
    prs.save(out)
    print(f"saved {out}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
